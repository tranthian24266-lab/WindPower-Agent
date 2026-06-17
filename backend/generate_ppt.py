import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Layouts
TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]

# Path to screenshots
FRONTEND_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "frontend")
SCREENSHOTS_DIR = os.path.join(FRONTEND_DIR, "test-reports", "screenshots")

def add_slide(title_text, bullets=None, image_path=None):
    slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
    
    # Set title
    title = slide.shapes.title
    title.text = title_text
    
    # Adjust title font
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.name = 'Microsoft YaHei'
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
    
    # Add bullets if present
    if bullets:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, text in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = text
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(18)
            p.level = 0
            
    # Resize and reposition text box if image is present
    if image_path and os.path.exists(image_path):
        if bullets:
            # Shift text box to the left
            body_shape = slide.shapes.placeholders[1]
            body_shape.width = Inches(4)
            body_shape.height = Inches(4.5)
            body_shape.left = Inches(0.5)
            body_shape.top = Inches(2)
            
            # Place image on the right
            slide.shapes.add_picture(image_path, Inches(4.8), Inches(2), width=Inches(4.8))
        else:
            # Center image if no bullets
            slide.shapes.add_picture(image_path, Inches(1.5), Inches(2), width=Inches(7))

# --- Slide 1: Title ---
title_slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]
title.text = "风电小模型库智能诊断平台"
subtitle.text = "全流程智能诊断闭环演示"
for p in title.text_frame.paragraphs: p.font.name = 'Microsoft YaHei'
for p in subtitle.text_frame.paragraphs: p.font.name = 'Microsoft YaHei'

# --- Slide 2 ---
add_slide(
    "1. 项目背景与目标",
    [
        "风电设备运维需要快速诊断、寿命预测和异常识别",
        "本平台目标：小模型库 + 智能体 + 知识库 + 自动报告",
        "解决问题：让用户不用直接操作模型代码，也能完成智能诊断"
    ]
)

# --- Slide 3 ---
add_slide(
    "2. 平台整体架构",
    [
        "前端：模型库、诊断、案例、报告、智能体、知识库",
        "后端：FastAPI 接口、模型调用、文件管理、案例、报告",
        "模型侧：故障诊断、剩余寿命预测、异常检测三个小模型",
        "智能体侧：结合诊断结果、知识库和大模型生成解释"
    ],
    os.path.join(SCREENSHOTS_DIR, "1_dashboard.png")
)

# --- Slide 4 ---
add_slide(
    "3. 核心功能一：模型库管理",
    [
        "展示不同任务类型的小模型",
        "支持模型版本、别名、验证状态和路由信息",
        "作用：统一管理诊断、RUL、异常检测模型"
    ],
    os.path.join(SCREENSHOTS_DIR, "2_models.png")
)

# --- Slide 5 ---
add_slide(
    "4. 核心功能二：上传数据并执行诊断",
    [
        "用户上传 .csv / .mat / .npy / .npz 文件",
        "选择任务类型",
        "系统自动调用对应小模型",
        "输出预测结果、置信度、风险等级等信息"
    ],
    os.path.join(SCREENSHOTS_DIR, "3_diagnosis.png")
)

# --- Slide 6 ---
add_slide(
    "5. 核心功能三：案例管理与结果查看",
    [
        "每次诊断自动生成案例",
        "支持查看历史案例",
        "支持按任务类型、风险等级筛选",
        "案例详情中展示模型结果和分析信息"
    ],
    os.path.join(SCREENSHOTS_DIR, "4_cases.png")
)

# --- Slide 7 ---
add_slide(
    "6. 核心功能四：报告生成",
    [
        "根据诊断结果生成分析报告",
        "支持基础报告和增强报告",
        "增强报告结合知识库、相似案例和大模型解释",
        "可用于辅助运维决策"
    ],
    os.path.join(SCREENSHOTS_DIR, "5_reports.png")
)

# --- Slide 8 ---
add_slide(
    "7. 核心功能五：智能体问答与知识库",
    [
        "用户可以围绕诊断案例继续提问",
        "智能体结合案例结果、模型信息、知识库内容回答",
        "知识库支持领域知识、模型说明、数据集说明检索",
        "把“模型输出”转化为“可理解的运维建议”"
    ],
    os.path.join(SCREENSHOTS_DIR, "6_chat.png")
)

# --- Slide 9 ---
add_slide(
    "8. 平台亮点与创新点",
    [
        "三类风电小模型统一接入",
        "前后端形成完整闭环",
        "诊断结果、知识库和大模型结合",
        "不只是模型预测，而是形成“诊断—解释—报告—问答”闭环",
        "具备扩展新模型、新知识、新设备档案的能力"
    ],
    os.path.join(SCREENSHOTS_DIR, "7_knowledge.png")
)

# --- Slide 10 ---
add_slide(
    "9. 当前完成度与后续计划",
    [
        "已完成：全流程功能打通、智能体问答、知识库、前端页面",
        "收尾优化：一键演示样例、界面细节、部署说明",
        "后续扩展：批量诊断、设备档案、健康趋势看板"
    ]
)

# --- Slide 11 ---
add_slide(
    "10. 总结",
    [
        "本平台已经形成较完整的风电智能诊断闭环",
        "可以支撑故障诊断、寿命预测、异常检测三类任务",
        "智能体和知识库极大增强了结果解释能力",
        "后续可继续向工程化运维平台扩展"
    ]
)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "风电智能诊断平台路演.pptx")
prs.save(output_path)
print(f"Presentation saved successfully at {output_path}")
